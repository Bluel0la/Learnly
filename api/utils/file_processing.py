import os
import io
import re
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
from fastapi import HTTPException


ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".pptx"}
MAX_FILE_SIZE_MB = 10


def extract_text_from_file(file_path: str) -> str:
    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension == ".docx":
        return extract_text_from_docx(file_path)
    elif extension == ".pptx":
        return extract_text_from_pptx(file_path)
    elif extension == ".txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError("Unsupported file type")


# PDF extraction using PyPDF2
def extract_text_from_pdf(file_path: str) -> str:
    text = []
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n".join(text)


# DOCX extraction
def extract_text_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])


# PPTX extraction
def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)


# TXT extraction
def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def clean_extracted_text(raw_text: str) -> str:
    # Removes multiple newlines and trims whitespace
    cleaned = re.sub(r"\n{2,}", "\n", raw_text)  # Collapse multiple newlines
    cleaned = re.sub(r"[ \t]+", " ", cleaned)  # Normalize whitespace
    return cleaned.strip()
